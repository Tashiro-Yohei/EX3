import sys
import streamlit as st
import streamlit.components.v1 as components
from google import genai
from google.genai import types
import pandas as pd
import json
import re
import time
from pptx import Presentation
import docx
import plotly.graph_objects as go
from io import BytesIO

try:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos
    _FPDF_AVAILABLE = True
except ImportError:
    _FPDF_AVAILABLE = False

# --- 1. モデル設定 & UI ---
st.set_page_config(page_title="AI Analyzer - Unified Score", layout="wide", initial_sidebar_state="expanded")

st.html("""
    <style>
        [data-testid="stSidebarHeader"] {
            padding-top: 0.5rem !important;
            padding-bottom: 0rem !important;
            min-height: auto !important;
        }
        [data-testid="stSidebarUserContent"] {
            padding-top: 0rem !important;
        }
        .block-container {
            padding-top: 3.5rem !important;
            padding-bottom: 2rem !important;
        }

        /* PDF出力・印刷用のカスタムCSS */
        @media print {
            [data-testid="stSidebar"] { display: none !important; }
            header[data-testid="stHeader"] { display: none !important; }
            .no-print { display: none !important; }
            iframe { display: none !important; }
            .block-container { padding-top: 0rem !important; padding-bottom: 0rem !important; max-width: 100% !important; }
            * {
                -webkit-print-color-adjust: exact !important;
                color-adjust: exact !important;
                print-color-adjust: exact !important;
            }

            /* A4印刷時の横並び継続とテキスト降下防止 */
            [data-testid="stHorizontalBlock"] {
                display: flex !important;
                flex-direction: row !important;
                flex-wrap: nowrap !important;
                align-items: flex-start !important;
            }
            [data-testid="column"] {
                min-width: 0 !important;
            }

            /* Plotly（グラフ）の強制表示 */
            .js-plotly-plot, .plotly, .plotly svg {
                max-width: 100% !important;
                height: auto !important;
            }
        }
    </style>
""")

# タイトルとエリア
st.html("""
    <div style="background-color:#f8f9fa; padding:20px; border-radius:10px; border-left:8px solid #0ea5e9; margin-bottom:25px;">
        <h2 style="margin:0; color:#1e293b;">🤖 AI Analyzer Pro - 統一スコア版</h2>
        <p style="margin:5px 0 0 0; color:#64748b; font-size:15px;">生成AIが出力したデータ内でのブランド言及が、企業理想とどれだけ合致しているかを分析します。スコア：低いほどヤバい。</p>
    </div>
""")

# セッション状態の初期化
if "ai_result" not in st.session_state:
    st.session_state.ai_result = None

# サイドバー：設定とエリア
with st.sidebar:
    st.markdown("### ⚙️ 解析設定")
    st.divider()

    api_key = st.text_input("🔑 Gemini API Key", type="password", help="Google AI Studioで発行されたAPIキーを入力してください")
    brand_name = st.text_input("🏷️ ブランド名", value="ディアナチュラ")
    brand_url = st.text_input("📍 公式サイトURL", value="https://www.dear-natura.com/")

    st.divider()
    st.markdown("### 📁 データのアップロード")
    uploaded_pptx = st.file_uploader("① ブランド戦略資料（PPTX）", type=["pptx"])
    uploaded_csv = st.file_uploader("② 生成AIでの言及データ(CSV)", type=["csv", "txt"])
    uploaded_docx = st.file_uploader("③ 生成AIの分析（DOCX）", type=["docx"])

    st.divider()
    if st.button("🔄 処理を中断 / 画面をリセット", type="secondary", use_container_width=True):
        st.session_state.ai_result = None
        st.rerun()

    st.divider()
    debug_mode = st.checkbox("🧠 デバッグ用ログを表示する", value=False)


# --- ユーティリティ関数 ---
def extract_text_from_pptx(file) -> str:
    try:
        prs = Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)[:3000]
    except Exception as e:
        return f"PPTXエラー: {e}"

def extract_text_from_docx(file) -> str:
    try:
        doc = docx.Document(file)
        text = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        return "\n".join(text)[:3000]
    except Exception as e:
        return f"DOCXエラー: {e}"

def load_csv_data(file):
    encodings = ["utf-8", "shift_jis", "cp932", "utf-8-sig", "iso-8859-1"]
    delimiters = [",", "\t", ";"]
    for encoding in encodings:
        for delimiter in delimiters:
            try:
                file.seek(0)
                df = pd.read_csv(file, encoding=encoding, sep=delimiter, engine="python", on_bad_lines="skip")
                if not df.empty and len(df.columns) >= 2:
                    return df
            except Exception:
                continue
    return None


# --- スコア定義表 ---
def score_to_level(score):
    """スコア（0-100）をレベルに変換。低いほどヤバい。"""
    score = int(score)
    if score >= 80:
        return {
            "level": "A",
            "emoji": "💚",
            "label": "達成済み（完璧）",
            "severity": 0
        }
    elif score >= 60:
        return {
            "level": "B",
            "emoji": "🟢",
            "label": "概ね達成",
            "severity": 1
        }
    elif score >= 40:
        return {
            "level": "C",
            "emoji": "🟡",
            "label": "課題あり",
            "severity": 2
        }
    elif score >= 20:
        return {
            "level": "D",
            "emoji": "🟠",
            "label": "大きな課題",
            "severity": 3
        }
    else:
        return {
            "level": "E",
            "emoji": "🔴",
            "label": "深刻な課題",
            "severity": 4
        }


# === PDF レポート生成関数 ===
def generate_pdf_report(res: dict, brand_name: str) -> bytes:
    """画面のすべての内容を含む詳細な PDF レポートを生成"""
    if not _FPDF_AVAILABLE:
        return None

    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(True, margin=12)
        pdf.set_margins(12, 12, 12)

        # フォント設定
        try:
            pdf.add_font("jp", "", "C:\\Windows\\Fonts\\YuGothR.ttc")
            pdf.add_font("jp", "B", "C:\\Windows\\Fonts\\YuGothB.ttc")
        except:
            pass

        pdf.set_font("jp", "B", 16)
        pdf.cell(0, 8, "AI Analyzer Pro - 統一スコア版", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 10)
        pdf.cell(0, 6, f"ブランド: {brand_name}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(4)

        story = res.get("diagnosis_story", {})
        unified_scores = res.get("unified_scores", {})
        unified_summary = res.get("unified_summary", "")
        unified_reasons = res.get("unified_reasons", {})
        scores = [unified_scores.get(k, 0) for k in ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]]
        overall_avg = int(sum(scores) / len(scores)) if scores else 0

        # === 総合評価 ===
        pdf.set_font("jp", "B", 12)
        pdf.cell(0, 7, "📊 総合評価", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 9)
        pdf.multi_cell(0, 4, f"総合スコア: {overall_avg}点\n{unified_summary}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(2)

        # === 現状診断 ===
        pdf.set_font("jp", "B", 12)
        pdf.cell(0, 7, "🎯 ① 現状の診断：企業理想とAIの言及のギャップ", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)

        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "✅ 正しく伝わっている強み", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, story.get("match", ""), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(1)

        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "💡 AIが見つけた予想外の強み", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, story.get("positive_gap", ""), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(1)

        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "⚠️ AIの言及が企業理想と乖離している点", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, story.get("negative_gap", ""), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(2)

        # === 統一スコア詳細 ===
        pdf.set_font("jp", "B", 12)
        pdf.cell(0, 7, "📈 ② 統一スコアによるAI言及合致度評価", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)

        score_labels = ["理念の合致度", "機能価値の合致度", "感情的価値の合致度", "安全性・評判", "シーン・モーメント"]
        score_keys = ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]

        for label, key in zip(score_labels, score_keys):
            score = unified_scores.get(key, 0)
            reason = unified_reasons.get(key, "")[:80]
            pdf.set_font("jp", "B", 9)
            pdf.cell(0, 4, f"{label}: {score}点", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_font("jp", "", 7)
            pdf.multi_cell(0, 3, f"  {reason}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(2)

        # === 対競業分析 ===
        comp_data = res.get("competitive_analysis", {})
        pdf.set_font("jp", "B", 12)
        pdf.cell(0, 7, "⚔️ ③ 対競業分析", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)

        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "ベンチマーク・競合企業との比較", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, comp_data.get("benchmark_competitors", ""), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "AI言及量の比較", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, comp_data.get("mention_volume_comparison", ""), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "優先順位の比較", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, comp_data.get("mention_order_comparison", ""), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "言及内容の比較", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, comp_data.get("mention_content_comparison", ""), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "戦略的アドバイス", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, comp_data.get("strategic_advice", ""), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(2)

        # === 改善施策 ===
        pdf.set_font("jp", "B", 12)
        pdf.cell(0, 7, "🎯 ④ 今後のマーケティング施策", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)

        topline = res.get("topline", "")
        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "戦略トップライン", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        pdf.multi_cell(0, 3, topline, new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        improvement_actions = res.get("improvement_actions", [])
        pdf.set_font("jp", "B", 9)
        pdf.cell(0, 4, "即座に取るべき5つのアクション", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)
        for i, action in enumerate(improvement_actions[:5], 1):
            pdf.multi_cell(0, 3, f"{i}. {action}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(2)

        # === 詳細乖離分析 ===
        pdf.add_page()
        pdf.set_font("jp", "B", 12)
        pdf.cell(0, 7, "🔍 ⑤ AI言及と企業理想の詳細な乖離分析", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)

        discrepancies = res.get("detailed_discrepancies", [])
        for i, disc in enumerate(discrepancies[:15], 1):
            if pdf.get_y() > 260:
                pdf.add_page()

            pdf.set_font("jp", "B", 9)
            pdf.cell(0, 4, f"乖離{i}: {disc.get('issue', '')[:60]}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

            pdf.set_font("jp", "", 8)
            impact = disc.get("impact", "")
            solution = disc.get("solution", "")
            pdf.multi_cell(0, 3, f"ビジネス影響: {impact[:70]}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.multi_cell(0, 3, f"解決策: {solution[:70]}\n", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # === トークデザイン案 ===
        pdf.add_page()
        pdf.set_font("jp", "B", 12)
        pdf.cell(0, 7, "💬 ⑥ トークデザイン開発（優先度別）", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("jp", "", 8)

        talk_designs = res.get("talk_designs", [])
        for i, td in enumerate(talk_designs[:10], 1):
            if pdf.get_y() > 250:
                pdf.add_page()

            pdf.set_font("jp", "B", 9)
            hook = td.get("talk_hook", "")[:50]
            roi = td.get("estimated_roi_score", 0)
            pdf.cell(0, 4, f"#{i} {hook} (ROI: {roi}%)", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

            pdf.set_font("jp", "", 8)
            impact = td.get("impact_score", 0)
            feasibility = td.get("feasibility_score", 0)
            pdf.cell(0, 3, f"  効果: {impact} | 実行性: {feasibility}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

            context = td.get("talk_context", "")[:80]
            asset = td.get("talk_asset", "")[:80]
            circulation = td.get("talk_circulation", "")[:80]

            pdf.multi_cell(0, 3, f"  コンテキスト: {context}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.multi_cell(0, 3, f"  アセット: {asset}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.multi_cell(0, 3, f"  サーキュレーション: {circulation}\n", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        return bytes(pdf.output())

    except Exception as e:
        st.error(f"PDF生成エラー: {e}")
        return None


def generate_level_interpretation(level_key, metric_name, is_ai=True):
    """スコアレベルの詳細解釈を生成"""
    interpretations = {
        "A": {
            "ai": f"{metric_name}について、生成AIの言及内容が企業理想と完全に合致している状態。",
            "action": "現状維持＋さらなる強化"
        },
        "B": {
            "ai": f"{metric_name}について、生成AIの言及内容が企業理想とほぼ合致している状態。",
            "action": "継続的な精度向上"
        },
        "C": {
            "ai": f"{metric_name}について、生成AIの言及があるが、企業理想との乖離がある状態。",
            "action": "学習データやプロンプト調整が必要"
        },
        "D": {
            "ai": f"{metric_name}について、生成AIの言及が限定的で、企業理想とのズレが大きい状態。",
            "action": "早期の改善施策が必要"
        },
        "E": {
            "ai": f"{metric_name}について、生成AIがほぼ言及していない、または完全に逆向きの状態。",
            "action": "根本的な改善が急務"
        }
    }

    return interpretations.get(level_key, {})


# --- 2. 解析ロジック ---
if st.button("📊 統一スコアで分析を実行", type="primary", use_container_width=True):
    if not api_key or not brand_url or not uploaded_pptx or not uploaded_csv or not uploaded_docx:
        st.error("左側のサイドバーで、APIキー・URL・3つのファイルがすべてセットアップされていません")
        st.stop()

    with st.spinner("AIデータを解析し、統一スコアベースのレポートを作成中..."):
        try:
            client = genai.Client(api_key=api_key)
            model_name = "gemini-2.5-flash"

            pptx_text = extract_text_from_pptx(uploaded_pptx)
            df_raw = load_csv_data(uploaded_csv)
            docx_text = extract_text_from_docx(uploaded_docx)

            if df_raw is None:
                st.error("CSVファイルの読み込みに失敗しました。")
                st.stop()

            csv_context = df_raw.head(35).to_csv(index=False)

            # Phase 1: オウンドメディアの企業理想辞書を抽出
            prompt_dict = f"""
            ブランド "{brand_name}" (公式URL: {brand_url}) の戦略資料からオウンドメディアが目指すブランド像を示す理想キーワードを5つずつ抽出してください。

            [戦略資料テキスト]
            {pptx_text}
            """

            dict_schema = {
                "type": "object",
                "properties": {
                    "core": {"type": "array", "items": {"type": "string"}},
                    "functional": {"type": "array", "items": {"type": "string"}},
                    "professional": {"type": "array", "items": {"type": "string"}}
                },
                "required": ["core", "functional", "professional"]
            }

            dictionary_data = None
            for attempt in range(3):
                try:
                    res_dict = client.models.generate_content(
                        model=model_name,
                        contents=[prompt_dict],
                        config=types.GenerateContentConfig(temperature=0.1, response_mime_type="application/json", response_schema=dict_schema)
                    )
                    dictionary_data = json.loads(res_dict.text)
                    break
                except Exception as e:
                    if "503" in str(e) and attempt < 2:
                        time.sleep(3)
                        continue
                    raise e

            if not dictionary_data:
                st.error("AIサーバーが混乱しています。少し時間を置いて再度試してください")
                st.stop()

            # Phase 2: 統一スコアベースの分析（修正版）
            response_schema = {
                "type": "object",
                "properties": {
                    "diagnosis_story": {
                        "type": "object",
                        "properties": {
                            "match": {"type": "string"},
                            "positive_gap": {"type": "string"},
                            "negative_gap": {"type": "string"}
                        },
                        "required": ["match", "positive_gap", "negative_gap"]
                    },
                    "topline": {"type": "string"},
                    "competitive_analysis": {
                        "type": "object",
                        "properties": {
                            "benchmark_competitors": {"type": "string"},
                            "mention_volume_comparison": {"type": "string"},
                            "mention_order_comparison": {"type": "string"},
                            "mention_content_comparison": {"type": "string"},
                            "strategic_advice": {"type": "string"}
                        },
                        "required": ["benchmark_competitors", "mention_volume_comparison", "mention_order_comparison", "mention_content_comparison", "strategic_advice"]
                    },
                    "improvement_actions": {"type": "array", "items": {"type": "string"}},
                    "detailed_discrepancies": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "issue": {"type": "string"},
                                "impact": {"type": "string"},
                                "solution": {"type": "string"}
                            },
                            "required": ["issue", "impact", "solution"]
                        }
                    },
                    "talk_designs": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "talk_hook": {"type": "string"},
                                "talk_context": {"type": "string"},
                                "talk_asset": {"type": "string"},
                                "talk_circulation": {"type": "string"},
                                "estimated_roi_score": {"type": "integer"},
                                "roi_reasoning": {"type": "string"},
                                "impact_score": {"type": "integer"},
                                "feasibility_score": {"type": "integer"}
                            },
                            "required": ["talk_hook", "talk_context", "talk_asset", "talk_circulation", "estimated_roi_score", "roi_reasoning", "impact_score", "feasibility_score"]
                        }
                    },
                    "unified_scores": {
                        "type": "object",
                        "properties": {
                            "brand_philosophy": {"type": "integer"},
                            "functional_value": {"type": "integer"},
                            "emotional_engagement": {"type": "integer"},
                            "safety_reputation": {"type": "integer"},
                            "usage_scene_moment": {"type": "integer"}
                        },
                        "required": ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]
                    },
                    "unified_summary": {"type": "string"},
                    "unified_reasons": {
                        "type": "object",
                        "properties": {
                            "brand_philosophy": {"type": "string"},
                            "functional_value": {"type": "string"},
                            "emotional_engagement": {"type": "string"},
                            "safety_reputation": {"type": "string"},
                            "usage_scene_moment": {"type": "string"}
                        },
                        "required": ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]
                    }
                },
                "required": ["diagnosis_story", "topline", "competitive_analysis", "improvement_actions", "detailed_discrepancies", "talk_designs", "unified_scores", "unified_summary", "unified_reasons"]
            }

            # ========== 修正版プロンプト（重要な変更） ==========
            prompt_analysis = f"""
            Analyze Generative AI output data for "{brand_name}" (Official URL: {brand_url}) against the owned media keywords.

            [企業の理想キーワード（Dictionary）]
            {json.dumps(dictionary_data, ensure_ascii=False)}

            [生成AIが出力したランキングデータ（CSV）]
            {csv_context}

            [生成AIの分析・説明（DOCX）]
            {docx_text}

            【修正版評価軸：AI言及合致度 0-100（低いほどヤバい）】

            従来：「企業理想をAIがどれだけ理解したか」
            修正後：「AIが出力したランキング・分析内でのブランド言及が、企業理想とどれだけ合致しているか」

            つまり：
            - CSVのランキング内でのブランド言及
            - DOCXの分析内でのブランド説明
            これらが、企業の Dictionary（理念・機能・感情・安全・シーン）とどれだけ合致しているかを測定

            TASK:

            1. "diagnosis_story": 3つのナラティブ（各200-250文字程度）を日本語で経営層向けに作成
               - "match": AIの言及内容が企業理想と正しく合致している部分
               - "positive_gap": AIの言及から発見できた予想外の強み
               - "negative_gap": AIの言及が企業理想と乖離している、または不足している部分

            2. "topline": 経営層向けの1行戦略

            3. "competitive_analysis":
               - "benchmark_competitors": 比較対象企業（定性的に1-2社を推奨、理由を説明）
               - "mention_volume_comparison": AI言及量の比較
               - "mention_order_comparison": AI言及の順位・優先度の比較
               - "mention_content_comparison": AI言及内容の比較
               - "strategic_advice": 戦略アドバイス

            4. "improvement_actions": 5つの改善施策

            5. "detailed_discrepancies": 最大10の具体的な乖離（AIの言及が企業理想とどう異なるか）
               - "issue": 乖離の詳細
               - "impact": ビジネスへの影響
               - "solution": 解決策

            6. "talk_designs": 10個のトークデザイン案（ROI優先度順）

            7. 【修正版】"unified_scores": AI言及の企業理想合致度を 0-100 で評価
               スコア計算方法：

               【理念の浸透度・合致度】
               - 0点：AIが言及しない or 完全に逆向き
               - 20点：少量の言及、大きなズレ
               - 40点：中程度の言及、部分的ズレ
               - 60点：多くの言及、概ね合致
               - 80点以上：十分な言及、ほぼ完全に合致

               同じ基準で5つの指標を評価してください：
               - brand_philosophy: 理念の言及合致度
               - functional_value: 機能価値の言及合致度
               - emotional_engagement: 感情的価値の言及合致度
               - safety_reputation: 安全性・評判の言及合致度
               - usage_scene_moment: シーン・モーメントの言及合致度

            8. "unified_summary": 5指標の平均スコアに基づく総評

            9. "unified_reasons": 各指標の詳細な評価理由

            Return JSON in Japanese.
            """
            # ========== プロンプト終了 ==========

            final_data = None
            for attempt in range(3):
                try:
                    res_analysis = client.models.generate_content(
                        model=model_name,
                        contents=[prompt_analysis],
                        config=types.GenerateContentConfig(temperature=0.1, response_mime_type="application/json", response_schema=response_schema)
                    )
                    final_data = json.loads(res_analysis.text)
                    break
                except Exception as e:
                    if "503" in str(e) and attempt < 2:
                        time.sleep(3)
                        continue
                    raise e

            if final_data:
                final_data["dictionary"] = dictionary_data
                st.session_state.ai_result = final_data
            else:
                st.error("AIサーバーが混乱しています。少し時間を置いて再度試してください")
                st.stop()

        except Exception as e:
            st.error(f"分析エラーが発生しました: {e}")
            st.stop()


# --- 3. 結果表示 UI ---
if st.session_state.ai_result:
    res = st.session_state.ai_result

    # PDF出力ボタンの配置 - 直接PDF ダウンロード
    if res:
        # PDF生成して直接ダウンロード
        pdf_data = generate_pdf_report(res, brand_name)

        col_pdf = st.columns([10, 1])[1]
        with col_pdf:
            st.download_button(
                "📄 PDFで出力",
                data=pdf_data or b"",
                file_name=f"{brand_name}_analysis.pdf",
                mime="application/pdf",
                disabled=pdf_data is None,
                use_container_width=True
            )

        # HTML レポート生成 - 詳細版（画面の全内容を含む）
        def generate_html_report():
            story = res.get("diagnosis_story", {})
            unified_scores = res.get("unified_scores", {})
            unified_summary = res.get("unified_summary", "")
            unified_reasons = res.get("unified_reasons", {})
            comp_data = res.get("competitive_analysis", {})
            topline = res.get("topline", "")
            improvement_actions = res.get("improvement_actions", [])
            discrepancies = res.get("detailed_discrepancies", [])
            talk_designs = res.get("talk_designs", [])

            # スコア計算
            scores = [unified_scores.get(k, 0) for k in ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]]
            overall_avg = int(sum(scores) / len(scores)) if scores else 0

            # 改善アクションのHTML生成
            actions_html = "".join([f"<li style='margin-bottom: 10px;'>{action}</li>" for action in improvement_actions])

            # 乖離分析のHTML生成
            discrepancies_html = ""
            for disc in discrepancies:
                discrepancies_html += f"""
                <div style="background: #fff9e6; padding: 15px; margin: 10px 0; border-left: 4px solid #ff9800; border-radius: 4px;">
                    <strong style="color: #ff6f00;">乖離: {disc.get('issue', '')}</strong>
                    <p style="margin: 10px 0 5px 0;"><strong>ビジネス影響:</strong> {disc.get('impact', '')}</p>
                    <p><strong>解決策:</strong> {disc.get('solution', '')}</p>
                </div>
                """

            # トークデザイン案のHTML生成
            designs_html = ""
            for i, td in enumerate(talk_designs[:10], 1):
                designs_html += f"""
                <div style="page-break-inside: avoid; background: #f5f5f5; padding: 15px; margin: 15px 0; border-left: 4px solid #2196f3; border-radius: 4px;">
                    <strong style="font-size: 16px; color: #1976d2;">#{i} {td.get('talk_hook', '')}</strong>
                    <p style="margin: 10px 0;"><strong>ROI:</strong> {td.get('estimated_roi_score', 0)}% | <strong>効果:</strong> {td.get('impact_score', 0)} | <strong>実行性:</strong> {td.get('feasibility_score', 0)}</p>
                    <p><strong>コンテキスト:</strong> {td.get('talk_context', '')}</p>
                    <p><strong>アセット:</strong> {td.get('talk_asset', '')}</p>
                    <p><strong>サーキュレーション:</strong> {td.get('talk_circulation', '')}</p>
                    <p><strong>ROI根拠:</strong> {td.get('roi_reasoning', '')}</p>
                </div>
                """

            html_content = f"""
            <!DOCTYPE html>
            <html lang="ja">
            <head>
                <meta charset="UTF-8">
                <title>AI Analyzer - 統一スコア版 レポート</title>
                <style>
                    * {{ margin: 0; padding: 0; box-sizing: border-box; }}
                    body {{ font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; line-height: 1.6; color: #333; background: #fff; }}
                    .container {{ max-width: 900px; margin: 0 auto; padding: 40px; }}
                    h1 {{ color: #1e293b; border-bottom: 3px solid #0ea5e9; padding-bottom: 15px; margin-bottom: 30px; font-size: 28px; }}
                    h2 {{ color: #0ea5e9; margin-top: 40px; margin-bottom: 20px; font-size: 22px; border-left: 4px solid #0ea5e9; padding-left: 15px; }}
                    h3 {{ color: #475569; margin-top: 20px; margin-bottom: 15px; font-size: 16px; }}
                    .summary-box {{ background: #e0f2fe; border-left: 5px solid #0ea5e9; padding: 20px; margin: 20px 0; border-radius: 4px; }}
                    .three-col {{ display: grid; grid-template-columns: 1fr 1fr 1fr; gap: 15px; margin: 20px 0; }}
                    .col-item {{ background: #f8f9fa; padding: 15px; border-radius: 4px; border-left: 4px solid #2563eb; }}
                    .col-item.positive {{ border-left-color: #28a745; background: #e6f4ea; }}
                    .col-item.negative {{ border-left-color: #dc3545; background: #fce8e6; }}
                    .score-grid {{ display: grid; grid-template-columns: repeat(5, 1fr); gap: 10px; margin: 20px 0; }}
                    .score-card {{ background: #f0f9ff; padding: 15px; border-radius: 8px; text-align: center; border: 1px solid #bae6fd; }}
                    .score-value {{ font-size: 32px; font-weight: bold; color: #0ea5e9; }}
                    .score-label {{ font-size: 12px; color: #475569; margin-top: 8px; }}
                    .overall-score {{ background: #e0f2fe; padding: 30px; text-align: center; border-radius: 8px; margin: 20px 0; border: 2px solid #0ea5e9; }}
                    .overall-value {{ font-size: 48px; font-weight: bold; color: #0ea5e9; }}
                    table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
                    th, td {{ border: 1px solid #e2e8f0; padding: 12px; text-align: left; }}
                    th {{ background: #0ea5e9; color: white; font-weight: bold; }}
                    tr:nth-child(even) {{ background: #f8fafc; }}
                    tr:hover {{ background: #f1f5f9; }}
                    .item {{ background: #f8f9fa; padding: 12px; margin: 10px 0; border-left: 4px solid #2563eb; border-radius: 4px; }}
                    .section {{ margin-bottom: 40px; }}
                    ul {{ margin-left: 20px; }}
                    li {{ margin-bottom: 8px; }}
                    p {{ margin-bottom: 10px; line-height: 1.8; }}
                    .page-break {{ page-break-after: always; }}
                    @media print {{ body {{ margin: 0; padding: 0; }} .container {{ padding: 20px; }} }}
                </style>
            </head>
            <body>
                <div class="container">
                    <h1>🤖 AI Analyzer Pro - 統一スコア版 完全レポート</h1>

                    <!-- 総合評価 -->
                    <div class="section">
                        <h2>📊 総合評価</h2>
                        <div class="overall-score">
                            <div style="font-size: 18px; color: #475569; margin-bottom: 10px;">AI言及合致度の総合スコア</div>
                            <div class="overall-value">{overall_avg}点</div>
                            <div style="font-size: 14px; color: #475569; margin-top: 10px;">{unified_summary}</div>
                        </div>
                    </div>

                    <!-- 現状の診断 -->
                    <div class="section">
                        <h2>🎯 ① 現状の診断：企業理想とAIの言及のギャップ</h2>
                        <div class="three-col">
                            <div class="col-item positive">
                                <strong>✅ 正しく伝わっている強み</strong>
                                <p style="margin-top: 10px; font-size: 13px;">{story.get("match", "分析データなし")}</p>
                            </div>
                            <div class="col-item">
                                <strong>💡 AIが見つけた予想外の強み</strong>
                                <p style="margin-top: 10px; font-size: 13px;">{story.get("positive_gap", "分析データなし")}</p>
                            </div>
                            <div class="col-item negative">
                                <strong>⚠️ 乖離している点</strong>
                                <p style="margin-top: 10px; font-size: 13px;">{story.get("negative_gap", "分析データなし")}</p>
                            </div>
                        </div>
                    </div>

                    <!-- 統一スコア評価 -->
                    <div class="section">
                        <h2>📈 ② 統一スコアによるAI言及合致度評価</h2>
                        <div class="score-grid">
                            <div class="score-card">
                                <div class="score-value">{unified_scores.get("brand_philosophy", 0)}</div>
                                <div class="score-label">理念の合致度</div>
                                <div style="font-size: 11px; color: #666; margin-top: 5px;">{unified_reasons.get("brand_philosophy", "")[:50]}...</div>
                            </div>
                            <div class="score-card">
                                <div class="score-value">{unified_scores.get("functional_value", 0)}</div>
                                <div class="score-label">機能価値の合致度</div>
                                <div style="font-size: 11px; color: #666; margin-top: 5px;">{unified_reasons.get("functional_value", "")[:50]}...</div>
                            </div>
                            <div class="score-card">
                                <div class="score-value">{unified_scores.get("emotional_engagement", 0)}</div>
                                <div class="score-label">感情的価値の合致度</div>
                                <div style="font-size: 11px; color: #666; margin-top: 5px;">{unified_reasons.get("emotional_engagement", "")[:50]}...</div>
                            </div>
                            <div class="score-card">
                                <div class="score-value">{unified_scores.get("safety_reputation", 0)}</div>
                                <div class="score-label">安全性・評判</div>
                                <div style="font-size: 11px; color: #666; margin-top: 5px;">{unified_reasons.get("safety_reputation", "")[:50]}...</div>
                            </div>
                            <div class="score-card">
                                <div class="score-value">{unified_scores.get("usage_scene_moment", 0)}</div>
                                <div class="score-label">シーン・モーメント</div>
                                <div style="font-size: 11px; color: #666; margin-top: 5px;">{unified_reasons.get("usage_scene_moment", "")[:50]}...</div>
                            </div>
                        </div>
                    </div>

                    <!-- 対競業分析 -->
                    <div class="section page-break">
                        <h2>⚔️ ③ 対競業分析</h2>
                        <div class="summary-box">
                            <strong>📌 ベンチマーク・競合企業との比較</strong>
                            <p style="margin-top: 10px;">{comp_data.get("benchmark_competitors", "データなし")}</p>
                        </div>
                        <h3>AI言及量の比較</h3>
                        <p>{comp_data.get("mention_volume_comparison", "データなし")}</p>
                        <h3>優先順位の比較</h3>
                        <p>{comp_data.get("mention_order_comparison", "データなし")}</p>
                        <h3>言及内容の比較</h3>
                        <p>{comp_data.get("mention_content_comparison", "データなし")}</p>
                        <div class="summary-box">
                            <strong>💡 戦略的アドバイス</strong>
                            <p style="margin-top: 10px;">{comp_data.get("strategic_advice", "データなし")}</p>
                        </div>
                    </div>

                    <!-- 改善施策 -->
                    <div class="section">
                        <h2>🎯 ④ 今後のマーケティング施策</h2>
                        <div class="summary-box">
                            <strong>📌 戦略トップライン</strong>
                            <p style="margin-top: 10px;">{topline}</p>
                        </div>
                        <h3>⚡ 即座に取るべき5つのアクション</h3>
                        <ul>{actions_html}</ul>
                    </div>

                    <!-- 乖離分析 -->
                    <div class="section page-break">
                        <h2>🔍 ⑤ AI言及と企業理想の詳細な乖離分析</h2>
                        {discrepancies_html}
                    </div>

                    <!-- トークデザイン案 -->
                    <div class="section page-break">
                        <h2>💬 ⑥ トークデザイン開発（優先順位別）</h2>
                        <p style="color: #666; font-size: 13px; margin-bottom: 20px;">企業理想に近づくための10のトーク案。ROI優先度順で掲載。</p>
                        {designs_html}
                    </div>

                    <!-- フッター -->
                    <div style="margin-top: 60px; padding-top: 20px; border-top: 1px solid #e2e8f0; text-align: center; color: #999; font-size: 12px;">
                        <p>このレポートは AI Analyzer で自動生成されました</p>
                        <p>生成日時: {pd.Timestamp.now().strftime('%Y年%m月%d日 %H:%M')}</p>
                    </div>
                </div>
            </body>
            </html>
            """
            return html_content

        html_data = generate_html_report()

        col1, col2, col3 = st.columns([8, 1, 1])
        with col2:
            st.download_button(
                "📄 HTML形式",
                data=html_data.encode('utf-8'),
                file_name="report.html",
                mime="text/html",
                use_container_width=True
            )
        with col3:
            st.info("💡 HTML をダウンロード後、ブラウザで開いて Ctrl+P で PDF に変換できます")
    else:
        st.info("分析結果がありません")

    if debug_mode and res:
        with st.expander("🧠 デバッグ用：全JSON出力"):
            st.json(res)

    # ==========================================
    # ① 現状の診断：企業理想とAIの言及のギャップ
    # ==========================================
    st.markdown("### 🎯 ① 現状の診断：企業理想とAIの言及のギャップ")
    st.caption("企業で発信したいメッセージがAIにどう言及されているか、その合致度と乖離点を分析します")

    story = res.get("diagnosis_story", {})

    col1, col2, col3 = st.columns(3)
    with col1:
        st.html("<div style='background-color:#e6f4ea; padding:12px; border-radius:5px; border-left:5px solid #28a745; font-weight:bold; color:#137333; margin-bottom:10px;'>✅ 正しく伝わっている強み</div>")
        st.write(story.get("match", "分析データなし"))

    with col2:
        st.html("<div style='background-color:#e8f0fe; padding:12px; border-radius:5px; border-left:5px solid #007bff; font-weight:bold; color:#1a73e8; margin-bottom:10px;'>💡 AIが見つけた予想外の強み</div>")
        st.write(story.get("positive_gap", "分析データなし"))

    with col3:
        st.html("<div style='background-color:#fce8e6; padding:12px; border-radius:5px; border-left:5px solid #dc3545; font-weight:bold; color:#c5221f; margin-bottom:10px;'>⚠️ AIの言及が企業理想と乖離している点</div>")
        st.write(story.get("negative_gap", "分析データなし"))

    st.divider()

    # ==========================================
    # ② 統一スコアによるAI言及合致度評価
    # ==========================================
    st.markdown("### 📊 ② 統一スコアによるAI言及合致度評価")
    st.caption("AIが出力したランキング・分析内でのブランド言及が、企業理想とどれだけ合致しているか。スコア：低いほどヤバい（0点に近い=企業理想から遠い）")

    unified_scores = res.get("unified_scores", {})
    unified_summary = res.get("unified_summary", "サマリーデータなし")
    unified_reasons = res.get("unified_reasons", {})

    categories = ['理念の\n合致度', '機能価値\nの合致度', '感情的価値\nの合致度', '安全性\n・評判', 'シーン\n・モーメント']
    categories_closed = categories + [categories[0]]

    keys = ["brand_philosophy", "functional_value", "emotional_engagement", "safety_reputation", "usage_scene_moment"]
    scores = [unified_scores.get(k, 0) for k in keys]
    scores_closed = scores + [scores[0]]

    overall_avg = int(sum(scores) / len(scores)) if scores else 0

    col_chart, col_summary = st.columns([1, 1.2])

    with col_chart:
        fig_radar = go.Figure()
        fig_radar.add_trace(go.Scatterpolar(
            r=scores_closed,
            theta=categories_closed,
            fill='toself',
            name='合致度スコア',
            line_color='#0ea5e9',
            fillcolor='rgba(14, 165, 233, 0.2)'
        ))
        fig_radar.update_layout(
            polar=dict(radialaxis=dict(visible=True, range=[0, 100], ticksuffix="点")),
            showlegend=False,
            margin=dict(l=40, r=40, t=30, b=30),
            height=350
        )
        st.plotly_chart(fig_radar, use_container_width=True, config={'staticPlot': True})

    with col_summary:
        overall_level = score_to_level(overall_avg)
        st.html(f"""
        <div style="margin-top: 20px; padding: 25px; background-color: #f8f9fa; border-left: 6px solid #0ea5e9; border-radius: 6px; box-shadow: 0 2px 4px rgba(0,0,0,0.05);">
            <div style="font-weight: bold; color: #0ea5e9; margin-bottom: 12px; font-size: 16px;">
                🎯 AI言及合致度の総合評価
            </div>
            <div style="font-size: 32px; font-weight: bold; color: #1e293b; margin-bottom: 8px;">
                {overall_avg}点 {overall_level['emoji']}
            </div>
            <div style="font-size: 14px; color: #475569; font-weight: bold; margin-bottom: 12px;">
                レベル{overall_level['level']}：{overall_level['label']}
            </div>
            <div style="font-size: 14px; color: #333; line-height: 1.8;">
                {unified_summary}
            </div>
        </div>
        """)

    st.markdown("<br>", unsafe_allow_html=True)

    # 詳細理由をカード形式で表示
    st.markdown("#### 📈 指標別の詳細評価")

    for title, key in zip(['理念の合致度', '機能価値の合致度', '感情的価値の合致度', '安全性・評判の合致度', 'シーン・モーメント適応度'], keys):
        score_val = unified_scores.get(key, 0)
        reason = unified_reasons.get(key, 'データなし')
        level_info = score_to_level(score_val)

        st.html(f"""
        <div style="border-left: 5px solid #0ea5e9; background-color: #f8f9fa; padding: 15px; margin-bottom: 12px; border-radius: 4px; box-shadow: 0 1px 2px rgba(0,0,0,0.05);">
            <div style="font-weight: bold; font-size: 16px; margin-bottom: 8px; color: #333;">
                📍 {title} <span style="color:#0ea5e9; font-size:24px; margin-left:10px; font-weight:bold;">{score_val}点</span> <span style="font-size:16px;">{level_info['emoji']}</span>
            </div>
            <div style="background-color: white; padding: 10px; border-radius: 4px; margin-bottom: 8px; border-left: 3px solid {['#dc3545','#fd7e14','#ffc107','#28a745','#20c997'][level_info['severity']]};">
                <div style="font-size: 12px; color: #475569; font-weight: bold; margin-bottom: 4px;">
                    レベル{level_info['level']}：{level_info['label']}
                </div>
                <div style="font-size: 13px; color: #333;">
                    {reason}
                </div>
            </div>
        </div>
        """)

    st.divider()

    # ==========================================
    # ③ 対競争分析
    # ==========================================
    st.markdown("### ⚔️ ③ 対競争分析")
    st.caption("AIの出力データから見える、競合との比較ポジション")

    comp_data = res.get("competitive_analysis", {})
    bench_comp = comp_data.get("benchmark_competitors", "データなし")
    vol_comp = comp_data.get("mention_volume_comparison", "データなし")
    order_comp = comp_data.get("mention_order_comparison", "データなし")
    content_comp = comp_data.get("mention_content_comparison", "データなし")

    st.html(f"""
    <div style="background-color: #fcfaff; border-left: 5px solid #8b5cf6; padding: 18px; border-radius: 4px; box-shadow: 0 1px 3px rgba(0,0,0,0.05); margin-bottom: 25px;">
        <div style="font-weight: bold; color: #7c3aed; margin-bottom: 8px; font-size: 16px;">🏆 ベンチマーク・競合企業との理由</div>
        <div style="font-size: 15px; color: #333; line-height: 1.6;">{bench_comp}</div>
    </div>

    <div style="margin-bottom: 20px;">
        <div style="border-left: 4px solid #475569; background-color: #f8f9fa; padding: 15px; margin-bottom: 10px; border-radius: 0 4px 4px 0;">
            <div style="font-weight: bold; color: #334155; font-size: 14px; margin-bottom: 4px;">📢 AI言及量の比較</div>
            <div style="font-size: 14px; color: #333; line-height: 1.5;">{vol_comp}</div>
        </div>
        <div style="border-left: 4px solid #475569; background-color: #f8f9fa; padding: 15px; margin-bottom: 10px; border-radius: 0 4px 4px 0;">
            <div style="font-weight: bold; color: #334155; font-size: 14px; margin-bottom: 4px;">📊 言及の優先順位の比較</div>
            <div style="font-size: 14px; color: #333; line-height: 1.5;">{order_comp}</div>
        </div>
        <div style="border-left: 4px solid #475569; background-color: #f8f9fa; padding: 15px; border-radius: 0 4px 4px 0;">
            <div style="font-weight: bold; color: #334155; font-size: 14px; margin-bottom: 4px;">💬 言及内容の比較</div>
            <div style="font-size: 14px; color: #333; line-height: 1.5;">{content_comp}</div>
        </div>
    </div>

    <div style="background-color: #e8f0fe; border: 1px solid #c2d7fa; padding: 15px; border-radius: 8px;">
        <div style="font-weight: bold; color: #1a73e8; margin-bottom: 8px;">💡 競争力の戦略的アドバイス</div>
        <div style="font-size: 15px; color: #333; line-height: 1.6;">{comp_data.get('strategic_advice', 'データなし')}</div>
    </div>
    """)

    st.divider()

    # ==========================================
    # ④ 今後のマーケティング施策案
    # ==========================================
    st.markdown("### 🎯 ④ 今後のマーケティング施策案")
    st.caption("AI言及の乖離を埋め、企業理想との合致度を高めるための5つのアクション")

    st.info(f"**📌 戦略方向：** {res.get('topline')}")

    st.markdown("#### ✅ 即座に取るべき5つのアクション")
    for i, action in enumerate(res.get("improvement_actions", []), 1):
        st.html(f"""
        <div style="background-color: #ffffff; border: 1px solid #e2e8f0; border-radius: 8px; padding: 15px; margin-bottom: 12px; display: flex; align-items: flex-start; box-shadow: 0 1px 3px rgba(0,0,0,0.05);">
            <div style="background-color: #0ea5e9; color: white; border-radius: 50%; min-width: 28px; height: 28px; display: flex; justify-content: center; align-items: center; font-weight: bold; margin-right: 15px; flex-shrink: 0;">
                {i}
            </div>
            <div style="font-size: 15px; color: #334155; line-height: 1.5; padding-top: 2px;">
                {action}
            </div>
        </div>
        """)

    st.divider()

    # ==========================================
    # ⑤ 詳細な乖離分析
    # ==========================================
    st.markdown("### 🔍 ⑤ AI言及と企業理想の詳細な乖離分析")
    st.caption("AI出力がどのような部分で企業理想と乖離しているか、ビジネス影響度とともに示します")

    discrepancies = res.get("detailed_discrepancies", [])
    if discrepancies:
        for item in discrepancies:
            issue = item.get("issue", "")
            impact = item.get("impact", "")
            solution = item.get("solution", "")

            st.html(f"""
            <div style="border: 1px solid #e2e8f0; border-radius: 8px; padding: 18px; margin-bottom: 15px; background-color: #ffffff; box-shadow: 0 1px 3px rgba(0,0,0,0.05);">
                <div style="font-weight: bold; color: #b91c1c; font-size: 16px; margin-bottom: 10px; display: flex; align-items: center;">
                    <span style="margin-right: 8px; font-size: 18px;">⚠️</span> 乖離: {issue}
                </div>
                <div style="font-size: 14px; color: #334155; margin-bottom: 12px; padding-left: 26px;">
                    <span style="font-weight: bold; color: #475569; display: block; margin-bottom: 4px;">📊 ビジネスへの影響:</span>
                    {impact}
                </div>
                <div style="font-size: 14px; color: #334155; padding-top: 12px; padding-left: 26px; border-top: 1px dashed #cbd5e1;">
                    <span style="font-weight: bold; color: #0284c7; display: block; margin-bottom: 4px;">💡 解決策：</span>
                    {solution}
                </div>
            </div>
            """)
    else:
        st.write("詳細な乖離は見つかりませんでした。")

    st.divider()

    # ==========================================
    # ⑥ トークデザイン案
    # ==========================================
    st.markdown("### 🎨 ⑥ トークデザイン開発（AI言及改善案）")
    st.caption("企業理想に近づくための、10個のトークデザイン・マーケティング案（ROI優先度順）")

    talk_designs = res.get("talk_designs", [])
    if talk_designs:
        for i, td in enumerate(talk_designs, 1):
            hook = td.get("talk_hook", "")
            context = td.get("talk_context", "")
            asset = td.get("talk_asset", "")
            circulation = td.get("talk_circulation", "")
            roi_score = td.get("estimated_roi_score", 0)
            roi_reasoning = td.get("roi_reasoning", "")
            impact_score = td.get("impact_score", 0)
            feasibility_score = td.get("feasibility_score", 0)

            rank_icon = "🥇" if i == 1 else "🥈" if i == 2 else "🥉" if i == 3 else ""

            st.html(f"""
            <div style="border: 1px solid #cbd5e1; border-radius: 8px; margin-bottom: 25px; background-color: #ffffff; box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);">
                <div style="background-color: #f8fafc; border-bottom: 1px solid #e2e8f0; padding: 15px 20px; border-radius: 8px 8px 0 0; display: flex; justify-content: space-between; align-items: center;">
                    <div style="font-weight: bold; font-size: 18px; color: #1e293b;">
                        {rank_icon} 優先度{i} トーク・フック: <span style="color: #0ea5e9;">「{hook}」</span>
                    </div>
                    <div style="background-color: #0ea5e9; color: white; padding: 5px 15px; border-radius: 20px; font-weight: bold; font-size: 14px;">
                        推定ROI: {roi_score}%
                    </div>
                </div>

                <div style="padding: 20px;">
                    <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 20px; margin-bottom: 20px;">
                        <div>
                            <div style="font-weight: bold; color: #475569; font-size: 13px; margin-bottom: 5px; text-transform: uppercase;">🎯 トーク・コンテキスト（誰が・どこで）</div>
                            <div style="font-size: 15px; color: #334155; line-height: 1.5;">{context}</div>
                        </div>
                        <div>
                            <div style="font-weight: bold; color: #475569; font-size: 13px; margin-bottom: 5px; text-transform: uppercase;">📎 トーク・アセット（証拠・事例）</div>
                            <div style="font-size: 15px; color: #334155; line-height: 1.5;">{asset}</div>
                        </div>
                    </div>

                    <div style="margin-bottom: 20px;">
                        <div style="font-weight: bold; color: #475569; font-size: 13px; margin-bottom: 5px; text-transform: uppercase;">📻 トーク・サーキュレーション（流通戦略）</div>
                        <div style="font-size: 15px; color: #334155; line-height: 1.5;">{circulation}</div>
                    </div>

                    <div style="border-top: 1px dashed #cbd5e1; padding-top: 15px; display: flex; justify-content: space-between; align-items: flex-start;">
                        <div style="flex: 1; padding-right: 20px;">
                            <div style="font-weight: bold; color: #64748b; font-size: 13px; margin-bottom: 5px;">📈 ROI根拠（言及改善 → 購買転換予測）</div>
                            <div style="font-size: 14px; color: #475569; line-height: 1.5;">{roi_reasoning}</div>
                        </div>
                        <div style="display: flex; gap: 15px; flex-shrink: 0; padding-top: 5px;">
                            <div style="background-color: #f8fafc; border: 1px solid #e2e8f0; padding: 8px 12px; border-radius: 6px; text-align: center;">
                                <div style="font-size: 11px; color: #64748b; font-weight: bold; margin-bottom: 2px;">効果</div>
                                <div style="font-size: 16px; color: #0f172a; font-weight: bold;">{impact_score}</div>
                            </div>
                            <div style="background-color: #f8fafc; border: 1px solid #e2e8f0; padding: 8px 12px; border-radius: 6px; text-align: center;">
                                <div style="font-size: 11px; color: #64748b; font-weight: bold; margin-bottom: 2px;">実行性</div>
                                <div style="font-size: 16px; color: #0f172a; font-weight: bold;">{feasibility_score}</div>
                            </div>
                        </div>
                    </div>
                </div>
            </div>
            """)

    st.divider()

    # 参考情報
    with st.expander("📌 参考：企業理想キーワード（Dictionary）"):
        d = res.get("dictionary", {})
        c1, c2, c3 = st.columns(3)
        with c1:
            st.markdown("**🎯 コア・根本的価値**")
            st.write(", ".join(d.get("core", [])) if d.get("core") else "抽出なし")
        with c2:
            st.markdown("**⚙️ 機能・効能**")
            st.write(", ".join(d.get("functional", [])) if d.get("functional") else "抽出なし")
        with c3:
            st.markdown("**🛡️ 信頼性・専門性**")
            st.write(", ".join(d.get("professional", [])) if d.get("professional") else "抽出なし")

else:
    st.html("""
        <div style="text-align:center; padding:100px 20px; color:#94a3b8;">
            <p style="font-size:40px; margin:0;">🤖</p>
            <h4 style="margin:10px 0 0 0; color:#64748b;">データがセットアップされていません</h4>
            <p style="font-size:14px; margin:5px 0 0 0;">左側のサイドバーにAPIキーとファイルをセットアップし、分析ボタンを押してください</p>
        </div>
    """)
